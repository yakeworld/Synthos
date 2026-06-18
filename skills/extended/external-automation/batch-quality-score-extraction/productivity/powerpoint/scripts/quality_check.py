#!/usr/bin/env python3
"""PPT quality review — overflow, overlap, content leak check.
Usage: python3 quality_check.py <file.pptx>
"""
import sys, re
from pptx import Presentation

if len(sys.argv) < 2:
    print("Usage: python3 quality_check.py <file.pptx>")
    sys.exit(1)

prs = Presentation(sys.argv[1])
SW = 13.333; SH = 7.5

issues = []
for pg_idx, slide in enumerate(prs.slides, 1):
    shapes_data = []
    for sh in slide.shapes:
        try:
            l = sh.left / 914400; t = sh.top / 914400
            w = sh.width / 914400; h = sh.height / 914400
            r = l + w; b = t + h
            txt = ""
            if sh.has_text_frame:
                txt = sh.text_frame.text[:60].replace('\n', ' ')
            shapes_data.append((l, t, r, b, w, h, txt, sh.shape_id))
        except: pass

    # 1. Slide overflow (ignore full-page backgrounds)
    for (l, t, r, b, w, h, txt, sid) in shapes_data:
        if r > SW + 0.01 and not (l < 0.3 and w > 12.5):
            issues.append(f"P{pg_idx}: RIGHT={r:.2f}\" > {SW}\" | {txt[:50]}")
        if b > SH + 0.01 and not (t < 0.3 and h > 7.0):
            issues.append(f"P{pg_idx}: BOTTOM={b:.2f}\" > {SH}\" | {txt[:50]}")

    # 2. Shape overlap (skip tiny/full-page shapes)
    for i in range(len(shapes_data)):
        for j in range(i+1, len(shapes_data)):
            l1, t1, r1, b1, w1, h1, txt1, sid1 = shapes_data[i]
            l2, t2, r2, b2, w2, h2, txt2, sid2 = shapes_data[j]
            if w1 > 12.5 or w2 > 12.5: continue  # full-page bg
            if w1 < 0.15 or w2 < 0.15: continue   # tiny
            if h1 < 0.15 or h2 < 0.15: continue
            # Skip if one contains the other (nested card pattern)
            if l1 <= l2 and t1 <= t2 and r1 >= r2 and b1 >= b2: continue
            if l2 <= l1 and t2 <= t1 and r2 >= r1 and b2 >= b1: continue
            ox = max(0, min(r1, r2) - max(l1, l2))
            oy = max(0, min(b1, b2) - max(t1, t2))
            if ox > 0.05 and oy > 0.05 and ox * oy > 0.5:
                issues.append(f"P{pg_idx}: OVERLAP {ox*oy:.1f}in² | [{txt1[:25]}] x [{txt2[:25]}]")

    # 3. Card content overflow — check text boxes inside cards
    cards = [(l,t,r,b,txt) for l,t,r,b,w,h,txt,sid in shapes_data if 3<(r-l)<9 and (b-t)>1.5]
    for cl,ct,cr,cb,ctxt in cards:
        for l,t,r,b,w,h,txt,sid in shapes_data:
            if l >= cl - 0.05 and t >= ct - 0.05:
                if r > cr + 0.05:
                    issues.append(f"P{pg_idx}: RIGHT LEAK | card=[{ctxt[:25]}] text=[{txt[:25]}] r={r:.2f}>{cr:.2f}")
                if b > cb + 0.05:
                    issues.append(f"P{pg_idx}: BOTTOM LEAK | card=[{ctxt[:25]}] text=[{txt[:25]}] b={b:.2f}>{cb:.2f}")

if issues:
    print(f"=== FOUND {len(issues)} ISSUES ===\n")
    by_page = {}
    for iss in issues:
        pg = int(re.search(r'P(\d+)', iss).group(1))
        by_page.setdefault(pg, []).append(iss)
    for pg in sorted(by_page):
        print(f"--- Page {pg} ({len(by_page[pg])} issues) ---")
        for iss in by_page[pg][:5]:
            print(f"  {iss}")
    sys.exit(1)
else:
    print("ALL CLEAN — no layout issues found")
    sys.exit(0)
