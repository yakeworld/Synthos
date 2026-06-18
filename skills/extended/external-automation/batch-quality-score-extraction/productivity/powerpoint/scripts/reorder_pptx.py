#!/usr/bin/env python3
"""调整PPT页面顺序 — 通过zip直接修改ppt/presentation.xml中的sldIdLst。

用法：python reorder_pptx.py <src.pptx> <dst.pptx> <new_order>
  new_order: 逗号分隔的原始索引，如 "0,1,2,3,4,5,12,13,6,7,8,9,10,11"

核心原理：python-pptx保存时使用prs.part._element中的XML（非prs.slides._sldIdLst）。
直接修改ZIP内部XML是唯一可靠方式。
"""

import zipfile
import re
import sys
from pptx import Presentation

def reorder_pptx(src, dst, order_indices):
    """重新排列pptx中的幻灯片顺序。"""
    # 读取原文件
    prs = Presentation(src)
    total = len(prs.slides)
    
    # 提取所有sldId
    with zipfile.ZipFile(src, 'r') as zin:
        pres_xml = zin.read('ppt/presentation.xml')
        all_sldIds = re.findall(r'<p:sldId id="(\d+)" r:id="(rId\d+)"/>', pres_xml)
        
        if len(all_sldIds) != total:
            print(f"WARNING: {len(all_sldIds)} sldIds but {total} slides")
        
        # 按新顺序排列
        new_order = [all_sldIds[i] for i in order_indices]
        
        # 重建sldIdLst
        new_block = '\n'.join(f'    <p:sldId id="{s}" r:id="{r}"/>' for s, r in new_order)
        replacement = '<p:sldIdLst>\n' + new_block + '\n  </p:sldIdLst>'
        new_xml = re.sub(
            r'(<p:sldIdLst>).*?(</p:sldIdLst>)',
            replacement,
            pres_xml.decode(),
            flags=re.DOTALL
        )
        
        # 写入新文件
        with zipfile.ZipFile(dst, 'w', zipfile.ZIP_DEFLATED) as zout:
            zout.writestr('ppt/presentation.xml', new_xml.encode())
            for m in zin.namelist():
                if m != 'ppt/presentation.xml':
                    zout.writestr(m, zin.read(m))
    
    print(f"已重排: {src} -> {dst}")
    print(f"顺序: {' -> '.join(str(i) for i in order_indices)}")

if __name__ == '__main__':
    if len(sys.argv) < 4:
        print(f"用法: {sys.argv[0]} <src.pptx> <dst.pptx> <order>")
        print(f"例: {sys.argv[0]} a.pptx b.pptx 0,1,2,12,13,3,4,5,6,7,8,9,10,11")
        sys.exit(1)
    
    src = sys.argv[1]
    dst = sys.argv[2]
    order = [int(x) for x in sys.argv[3].split(',')]
    reorder_pptx(src, dst, order)
