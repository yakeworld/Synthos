#!/usr/bin/env python3
"""
将 PowerPoint（.pptx）按页导出为 Markdown，并抽取幻灯片中的嵌入图片，便于 Step 2 扫描与 Agent Read。

依赖 python-pptx（见仓库根目录 requirements.txt）。

用法:
  python pptx_to_md.py --input review.pptx --output outputs/case/review.md
  python pptx_to_md.py -i a.pptx -o b/out.md --media-dir b/slide_images

默认图片目录：与输出 .md 同级的「{md 文件名}_media/」。
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path


def _require_pptx():
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        print(
            "缺少依赖 python-pptx。请在技能根目录执行: pip install -r requirements.txt",
            file=sys.stderr,
        )
        sys.exit(1)
    return Presentation, MSO_SHAPE_TYPE


def _walk_shapes(shapes, MSO_SHAPE_TYPE):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes, MSO_SHAPE_TYPE)
        else:
            yield shape


def _shape_text(shape) -> str:
    if getattr(shape, "has_text_frame", False):
        t = (shape.text_frame.text or "").strip()
        return t
    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            cells = []
            for cell in row.cells:
                cells.append((cell.text or "").strip().replace("\n", " "))
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            return "\n".join(rows)
    return ""


def _rel_media_path(out_file: Path, media_file: Path) -> str:
    try:
        return media_file.relative_to(out_file.parent).as_posix()
    except ValueError:
        return media_file.as_posix()


def _run(input_pptx: Path, output_md: Path, media_dir: Path | None) -> int:
    Presentation, MSO_SHAPE_TYPE = _require_pptx()

    if not input_pptx.is_file():
        print(f"输入文件不存在: {input_pptx}", file=sys.stderr)
        return 2
    suf = input_pptx.suffix.lower()
    if suf not in (".pptx", ".ppsx"):
        print("警告: 期望 .pptx / .ppsx（OOXML）；旧版 .ppt 不支持。", file=sys.stderr)

    output_md = output_md.resolve()
    output_md.parent.mkdir(parents=True, exist_ok=True)

    if media_dir is None:
        media_dir = output_md.parent / f"{output_md.stem}_media"
    else:
        media_dir = media_dir.resolve()
    media_dir.mkdir(parents=True, exist_ok=True)

    try:
        prs = Presentation(str(input_pptx))
    except Exception as e:
        print(f"无法打开演示文稿: {e}", file=sys.stderr)
        return 3

    lines: list[str] = [
        f"<!-- 由 pptx_to_md.py 自 {input_pptx.name} 转换，勿手改本行元信息 -->\n"
    ]
    img_counter = [0]

    for sn, slide in enumerate(prs.slides, start=1):
        lines.append(f"\n## 第 {sn} 页\n")

        for shape in _walk_shapes(slide.shapes, MSO_SHAPE_TYPE):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    ext = (img.ext or "png").lower()
                    if ext == "jpeg":
                        ext = "jpg"
                    img_counter[0] += 1
                    fname = f"slide{sn:02d}_img{img_counter[0]:04d}.{ext}"
                    out_img = media_dir / fname
                    out_img.write_bytes(img.blob)
                    rel = _rel_media_path(output_md, out_img)
                    lines.append(f"\n![]({rel})\n")
                except Exception as e:
                    print(f"警告: 第 {sn} 页抽取图片失败: {e}", file=sys.stderr)
                continue

            block = _shape_text(shape)
            if block:
                lines.append(block)
                lines.append("\n\n")

        try:
            nf = slide.notes_slide.notes_text_frame
            note_txt = (nf.text or "").strip() if nf is not None else ""
            if note_txt:
                lines.append("\n**备注**：\n\n")
                lines.append(note_txt)
                lines.append("\n\n")
        except (AttributeError, ValueError):
            pass

    body = "".join(lines).rstrip() + "\n"
    output_md.write_text(body, encoding="utf-8")

    print(f"已写入: {output_md}")
    print(f"图片目录: {media_dir}")
    return 0


def main() -> int:
    p = argparse.ArgumentParser(description="PowerPoint (.pptx) → Markdown + 抽取图片")
    p.add_argument("-i", "--input", required=True, type=Path, help="输入 .pptx / .ppsx 路径")
    p.add_argument("-o", "--output", required=True, type=Path, help="输出 .md 路径")
    p.add_argument(
        "--media-dir",
        type=Path,
        default=None,
        help="图片输出目录（默认：与 .md 同级的 {md 主名}_media）",
    )
    args = p.parse_args()
    return _run(args.input, args.output, args.media_dir)


if __name__ == "__main__":
    raise SystemExit(main())
