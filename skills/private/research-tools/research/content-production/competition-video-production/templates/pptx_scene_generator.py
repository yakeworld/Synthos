#!/usr/bin/env python3
"""
PPTX scene generator for competition/demo videos.
User context: Wenzhou Hospital / Synthos project.
Dark theme with accent cards, data visualization, DAG flows.

Usage: Adapt the scene creation functions for your own content.
Each scene should be ~40-50 seconds of video time.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette: Tech Dark ───
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG = RGBColor(0x1E, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x84)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_80 = RGBColor(0xCC, 0xCC, 0xCC)
WHITE_50 = RGBColor(0x99, 0x99, 0x99)

OUTPUT_DIR = "/media/yakeworld/sda2/Synthos/docs/video_pptx"
os.makedirs(OUTPUT_DIR, exist_ok=True)


def add_background(slide, color=BG_DARK):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, fill_color=None,
              line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
        if line_width:
            shape.line.width = line_width
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox


def add_card(slide, left, top, width, height, title, body="",
             title_color=ACCENT_BLUE, body_color=WHITE_80,
             card_bg=CARD_BG, accent_bar=True):
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                     fill_color=card_bg, line_color=RGBColor(0x2A, 0x3A, 0x5C),
                     line_width=Pt(1))
    if accent_bar:
        bar = add_shape(slide, MSO_SHAPE.RECTANGLE, left, top,
                        Inches(0.06), height,
                        fill_color=title_color if isinstance(title_color, RGBColor) else title_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1),
                 width - Inches(0.3), Inches(0.4),
                 title, font_size=14, font_color=title_color, bold=True)
    if body:
        add_text_box(slide, left + Inches(0.2), top + Inches(0.45),
                     width - Inches(0.3), height - Inches(0.55),
                     body, font_size=11, font_color=body_color, line_spacing=1.3)
    return card


def add_icon_circle(slide, left, top, size, color, text="", font_size=16):
    circle = add_shape(slide, MSO_SHAPE.OVAL, left, top, size, size,
                       fill_color=color)
    if text:
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    return circle


def add_arrow_right(slide, left, top, width, height, color=ACCENT_BLUE, text=""):
    arrow = add_shape(slide, MSO_SHAPE.RIGHT_ARROW, left, top, width, height,
                      fill_color=color)
    if text:
        tf = arrow.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(9)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    return arrow


# ─── Template Scene Creator ───
def create_template_scene(prs, scene_number, title, subtitle, cards_data):
    """
    Create a scene slide with title, subtitle, and card grid.
    
    Args:
        prs: Presentation object
        scene_number: Scene number for reference
        title: Main title text
        subtitle: Subtitle or description
        cards_data: List of (title, body, accent_color) tuples
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_background(slide, BG_DARK)
    
    # Title
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
                 title, font_size=22, font_color=WHITE, bold=True)
    
    # Subtitle
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.5),
                     subtitle, font_size=16, font_color=WHITE_80)
    
    # Cards
    card_w = Inches(3.0)
    card_h = Inches(1.2)
    gap = Inches(0.3)
    y = Inches(1.5)
    
    for i, (card_title, card_body, card_color) in enumerate(cards_data):
        row = i // 3
        col = i % 3
        x = Inches(0.5) + col * (card_w + gap)
        add_card(slide, x, y + row * (card_h + gap),
                 card_w, card_h, card_title, card_body,
                 title_color=card_color, body_color=WHITE_80,
                 accent_bar=card_color)


def create_dag_flow(prs, title, nodes, arrow_color=ACCENT_BLUE):
    """
    Create a directed acyclic graph flow diagram.
    
    Args:
        nodes: List of (title, description, color) tuples
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG_DARK)
    
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
                 title, font_size=22, font_color=WHITE, bold=True)
    
    node_w = Inches(2.5)
    node_h = Inches(1.5)
    gap = Inches(0.4)
    y = Inches(1.2)
    
    for i, (name, desc, color) in enumerate(nodes):
        x = Inches(0.5) + i * (node_w + gap)
        
        # Node card
        node = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                         x, y, node_w, node_h,
                         fill_color=CARD_BG, line_color=color,
                         line_width=Pt(2))
        
        # Number badge
        add_icon_circle(slide, x + Inches(0.1), y + Inches(0.1),
                        Inches(0.4), color, str(i + 1), 14)
        
        # Title
        add_text_box(slide, x + Inches(0.55), y + Inches(0.15),
                     node_w - Inches(0.7), Inches(0.4),
                     name, font_size=13, font_color=color, bold=True)
        
        # Description
        add_text_box(slide, x + Inches(0.55), y + Inches(0.55),
                     node_w - Inches(0.7), Inches(0.8),
                     desc, font_size=11, font_color=WHITE_80, line_spacing=1.3)
        
        # Arrow to next
        if i < len(nodes) - 1:
            add_arrow_right(slide, x + node_w, y + Inches(0.6),
                            gap, Inches(0.35), arrow_color)


def create_comparison(prs, title, left_items, right_items, left_color, right_color):
    """Create before/after comparison layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG_DARK)
    
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
                 title, font_size=22, font_color=WHITE, bold=True)
    
    # Left column (before)
    add_card(slide, Inches(0.5), Inches(1.0), Inches(5.8), Inches(3.0),
             "Before", left_items,
             title_color=left_color, body_color=WHITE_50,
             accent_bar=left_color, card_bg=RGBColor(0x1A, 0x20, 0x30))
    
    # Right column (after)
    add_card(slide, Inches(7), Inches(1.0), Inches(5.8), Inches(3.0),
             "After", right_items,
             title_color=right_color, body_color=WHITE_80,
             accent_bar=right_color, card_bg=RGBColor(0x1A, 0x25, 0x40))
    
    # Center arrow
    add_arrow_right(slide, Inches(5.9), Inches(2.2), Inches(0.8), Inches(0.5),
                    right_color)


# ─── Example: Open/Title Slide ───
def create_title_slide(prs, logo_text, subtitle_text, tagline=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG_DARK)
    
    # Logo area
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(3.5), Inches(1.8), Inches(6.3), Inches(1.5),
              fill_color=RGBColor(0x1A, 0x25, 0x44),
              line_color=ACCENT_BLUE, line_width=Pt(2))
    
    add_text_box(slide, Inches(3.8), Inches(1.95), Inches(5.7), Inches(0.7),
                 logo_text, font_size=48, font_color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(3.5), Inches(2.65), Inches(6.3), Inches(0.4),
                 subtitle_text, font_size=22, font_color=ACCENT_BLUE,
                 alignment=PP_ALIGN.CENTER)
    
    if tagline:
        add_text_box(slide, Inches(2), Inches(4.0), Inches(9.3), Inches(0.5),
                     tagline, font_size=18, font_color=ACCENT_CYAN,
                     alignment=PP_ALIGN.CENTER)